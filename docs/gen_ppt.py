
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ORANGE = RGBColor(0xFF, 0x6D, 0x00)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xD3, 0x2F, 0x2F)
BLUE = RGBColor(0x15, 0x65, 0xC0)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_ORANGE = RGBColor(0xFF, 0xE0, 0xB2)
BG = RGBColor(0xFF, 0xFB, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.2, fill_color=ORANGE)
    add_text(slide, title, 0.4, 0.1, 10, 0.7, font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.75, 10, 0.4, font_size=16, color=LIGHT_ORANGE)

def add_bullet(slide, items, left, top, width, font_size=18, color=DARK):
    y = top
    for item in items:
        dot = slide.shapes.add_shape(9, Inches(left), Inches(y + 0.1), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE
        dot.line.fill.background()
        add_text(slide, item, left + 0.22, y, width - 0.22, 0.42, font_size=font_size, color=color)
        y += 0.46
    return y

def page_num(slide, n, total=10):
    add_text(slide, f"{n} / {total}", 12.3, 7.1, 0.9, 0.3, font_size=12, color=GRAY, align=PP_ALIGN.RIGHT)

# ===== Slide 1: 封面 =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=ORANGE)
# 装饰圆
for r, c in [(6.0, RGBColor(0xFF,0x8F,0x00)), (4.0, RGBColor(0xFF,0xA0,0x30)), (2.5, RGBColor(0xFF,0xB0,0x50))]:
    shape = s.shapes.add_shape(9, Inches(9.0 - r/2), Inches(4.5 - r/2), Inches(r), Inches(r))
    shape.fill.solid(); shape.fill.fore_color.rgb = c; shape.line.fill.background()
add_text(s, "信息服务APP的适老化设计", 0.8, 1.5, 9.0, 1.0, font_size=40, bold=True, color=WHITE)
add_text(s, "与多模态交互", 0.8, 2.5, 9.0, 0.9, font_size=40, bold=True, color=WHITE)
add_text(s, "核心系统：小浙 — 适老化浙里办长辈版", 0.8, 3.6, 9.0, 0.55, font_size=22, color=LIGHT_ORANGE)
add_text(s, "语音对话 · 代理填表 · 人脸认证 · 政务场景", 0.8, 4.2, 9.0, 0.5, font_size=18, color=LIGHT_ORANGE)
add_rect(s, 0.8, 5.5, 2.5, 0.05, fill_color=WHITE)
add_text(s, "成果展示答辩", 0.8, 5.65, 6.0, 0.5, font_size=16, color=WHITE)

# ===== Slide 2: 成果全景 =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=BG)
add_title_bar(s, "成果全景", "信息服务APP的适老化设计与多模态交互")
add_rect(s, 0.4, 1.35, 5.5, 5.75, fill_color=WHITE, line_color=RGBColor(0xE0,0xE0,0xE0), line_width=1)
add_text(s, "项目成果一览", 0.6, 1.45, 5.0, 0.45, font_size=20, bold=True, color=ORANGE)
bullets = [
    "可运行原型：17个页面，4个核心政务场景",
    "核心系统：「小浙」适老化浙里办长辈版",
    "核心能力：语音对话→理解意图→代填执行",
    "4场景：登录·医保缴费·社保查询·公告浏览",
    "手机浏览器直接可用，无需安装",
    "前端 Flutter Web + 后端 FastAPI + DeepSeek-V3",
    "讯飞 ASR/TTS + MediaPipe 本地人脸检测",
]
add_bullet(s, bullets, 0.6, 2.0, 5.1, font_size=16)
add_text(s, "界面全家福（截图区）", 6.3, 1.4, 6.6, 0.4, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
positions = [(6.3, 1.85), (9.6, 1.85), (6.3, 4.5), (9.6, 4.5)]
labels = ["长辈版首页", "代理面板弹出", "授权卡片", "查询结果大字"]
for i, (lx, ly) in enumerate(positions):
    add_rect(s, lx, ly, 3.0, 2.4, fill_color=RGBColor(0xE8,0xE8,0xE8), line_color=GRAY, line_width=0.5)
    add_text(s, labels[i], lx, ly + 0.95, 3.0, 0.4, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "[ 可替换为截图 ]", lx, ly + 1.35, 3.0, 0.4, font_size=11, color=RGBColor(0xBD,0xBD,0xBD), align=PP_ALIGN.CENTER)
page_num(s, 2)

# ===== Slide 3: 核心创新 =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=BG)
add_title_bar(s, "核心创新：受控响应型代理框架", "三原则固化为架构约束，从对话约定升级为可验证工程性质")
innovations = [
    ("原则一\n确定性隔离", ORANGE, [
        "高风险操作永不作为 AI 工具注册",
        "代码层物理隔离，提示词注入无效",
        "政务数据不离开合规环境",
    ]),
    ("原则二\n最小授权", GREEN, [
        "L1/L2/L3 三级分级授权矩阵",
        "权限随 WebSocket 会话消亡",
        "每步操作需用户显式确认",
    ]),
    ("原则三\n透明可解释", BLUE, [
        "Agent 每步操作实时播报",
        "草稿箱可随时查看/撤销/续填",
        "「代理在做什么」始终可见",
    ]),
]
cw = 3.9
for i, (title, color, items) in enumerate(innovations):
    cx = 0.4 + i * (cw + 0.3)
    add_rect(s, cx, 1.4, cw, 0.55, fill_color=color)
    add_text(s, title, cx + 0.1, 1.45, cw - 0.2, 0.45, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, cx, 1.95, cw, 3.8, fill_color=WHITE, line_color=color, line_width=1)
    y = 2.05
    for item in items:
        dot = s.shapes.add_shape(9, Inches(cx+0.15), Inches(y+0.12), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
        add_text(s, item, cx+0.35, y, cw-0.45, 0.65, font_size=15, color=DARK)
        y += 0.85
# 底部总结
add_rect(s, 0.4, 6.0, 12.5, 0.55, fill_color=ORANGE)
add_text(s, "信任保障 = 架构约束，而非对话约定 | 协助而非替代，透明而非黑箱", 0.6, 6.08, 12.1, 0.4, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
page_num(s, 3)

# ===== Slide 4: 多模态交互设计 =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=BG)
add_title_bar(s, "多模态交互设计", "语音·视觉·触控三通道冗余互补，任一退化不影响使用")
modalities = [
    ("语音通道", ORANGE, "讯飞 ASR/TTS", [
        "实时语音识别，流式转写",
        "自然语言理解意图",
        "TTS 播报操作结果",
        "MicButton 悬浮按钮随时可唤起",
    ]),
    ("视觉通道", BLUE, "MediaPipe 本地", [
        "人脸检测完全本地，零网络依赖",
        "刷脸登录无需额外 SDK",
        "大字模式，对比度符合 WCAG AA",
        "AgentElementRegistry 高亮定位",
    ]),
    ("触控通道", GREEN, "Flutter手势", [
        "最小点击区域 48×48 dp",
        "单步确认，防误触设计",
        "授权卡片滑动确认",
        "草稿箱断点续填",
    ]),
]
cw = 3.9
for i, (title, color, tech, items) in enumerate(modalities):
    cx = 0.4 + i * (cw + 0.3)
    add_rect(s, cx, 1.4, cw, 0.5, fill_color=color)
    add_text(s, title, cx+0.1, 1.45, cw-0.2, 0.4, font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, cx, 1.9, cw, 0.38, fill_color=RGBColor(0xEE,0xEE,0xEE))
    add_text(s, tech, cx+0.1, 1.93, cw-0.2, 0.32, font_size=13, color=GRAY, align=PP_ALIGN.CENTER)
    add_rect(s, cx, 2.28, cw, 3.5, fill_color=WHITE, line_color=color, line_width=1)
    y = 2.38
    for item in items:
        dot = s.shapes.add_shape(9, Inches(cx+0.15), Inches(y+0.12), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
        add_text(s, item, cx+0.35, y, cw-0.45, 0.6, font_size=15, color=DARK)
        y += 0.72
add_rect(s, 0.4, 6.0, 12.5, 0.55, fill_color=RGBColor(0xEE,0xEE,0xEE))
add_text(s, "冗余互补设计：感官通道退化（听力/视力下降）时，其余通道自动补偿，确保老年用户始终可用", 0.6, 6.08, 12.1, 0.4, font_size=14, color=DARK, align=PP_ALIGN.CENTER)
page_num(s, 4)

# ===== Slide 5: DEMO演示流程 =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=BG)
add_title_bar(s, "DEMO 演示：医保缴费全流程", "扫码访问 → 语音唤起 → 代理填表 → 用户确认 → 完成")
steps = [
    ("Step 1", "打开长辈版首页", "扫码或访问 localhost:3080\n大字界面，橙色「找小浙帮忙」按钮"),
    ("Step 2", "语音唤起代理", "点击麦克风按钮\n说：「帮我缴医保」"),
    ("Step 3", "意图理解", "DeepSeek-V3 解析意图\n匹配「医保缴费」场景"),
    ("Step 4", "代理面板弹出", "展示操作步骤预览\n用户点击「授权」"),
    ("Step 5", "代理自动填表", "逐字段高亮填写\n每步实时播报"),
    ("Step 6", "用户最终确认", "展示草稿，用户核对\n点击「提交」完成"),
]
sw = 2.0
for i, (step, title, desc) in enumerate(steps):
    sx = 0.35 + i * (sw + 0.05)
    add_rect(s, sx, 1.4, sw, 0.45, fill_color=ORANGE)
    add_text(s, step, sx+0.05, 1.44, sw-0.1, 0.37, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, sx, 1.85, sw, 4.6, fill_color=WHITE, line_color=RGBColor(0xE0,0xE0,0xE0), line_width=1)
    add_text(s, title, sx+0.1, 1.95, sw-0.2, 0.45, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_rect(s, sx+0.2, 2.45, sw-0.4, 0.03, fill_color=ORANGE)
    add_text(s, desc, sx+0.1, 2.55, sw-0.2, 2.5, font_size=13, color=GRAY, align=PP_ALIGN.CENTER)
    if i < 5:
        add_text(s, "→", sx + sw + 0.01, 2.8, 0.08, 0.4, font_size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_rect(s, 0.35, 6.55, 12.6, 0.5, fill_color=RGBColor(0xFF,0xEE,0xCC))
add_text(s, "演示地址：http://localhost:3080  |  建议投屏手机展示效果", 0.5, 6.6, 12.3, 0.38, font_size=14, color=ORANGE, align=PP_ALIGN.CENTER)
page_num(s, 5)

# ===== Slide 6: 适老化设计细节 =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=BG)
add_title_bar(s, "适老化设计细节", "以老年用户为中心的每一个细节决策")
details = [
    ("视觉适老", [
        "最小字号 18sp，标题 24-32sp",
        "主色橙色 #FF6D00，高对比度",
        "WCAG AA 对比度标准",
        "关键操作按钮 ≥ 56dp 高度",
    ]),
    ("认知减负", [
        "单屏单任务，去除多余信息",
        "步骤线性引导，不分叉",
        "错误提示用大字+语音双播报",
        "「后退」随时可用，无惩罚机制",
    ]),
    ("信任建立", [
        "代理操作全程可见、可暂停",
        "草稿箱：任何时候可以中断续填",
        "授权卡片明确展示将做什么",
        "「不授权」始终是默认安全态",
    ]),
    ("运动辅助", [
        "大触控区域防误触",
        "语音替代复杂文字输入",
        "刷脸替代密码记忆",
        "震动反馈确认关键操作",
    ]),
]
cw = 2.95
for i, (title, items) in enumerate(details):
    cx = 0.35 + i * (cw + 0.25)
    add_rect(s, cx, 1.4, cw, 0.45, fill_color=ORANGE)
    add_text(s, title, cx+0.1, 1.44, cw-0.2, 0.37, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, cx, 1.85, cw, 4.6, fill_color=WHITE, line_color=RGBColor(0xE0,0xE0,0xE0), line_width=1)
    y = 1.98
    for item in items:
        dot = s.shapes.add_shape(9, Inches(cx+0.15), Inches(y+0.12), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = ORANGE; dot.line.fill.background()
        add_text(s, item, cx+0.32, y, cw-0.42, 0.6, font_size=14, color=DARK)
        y += 0.72
page_num(s, 6)

# ===== Slide 7: 4大政务场景 =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=BG)
add_title_bar(s, "4大政务场景覆盖", "从登录到查询，完整的老年用户政务体验链路")
scenarios = [
    ("场景一\n登录认证", ORANGE, [
        "验证码登录（大字键盘）",
        "人脸识别登录（MediaPipe）",
        "记住登录态，减少重复操作",
    ]),
    ("场景二\n医保缴费", RGBColor(0xE6,0x3B,0x00), [
        "语音说出需求自动跳转",
        "代理逐字段填写表单",
        "金额确认卡片防误操作",
    ]),
    ("场景三\n社保查询", GREEN, [
        "余额/记录大字展示",
        "语音播报查询结果",
        "历史记录时间轴呈现",
    ]),
    ("场景四\n政务公告", BLUE, [
        "字号动态调节（A-/A+）",
        "朗读按钮全文语音播报",
        "重要通知红点角标提醒",
    ]),
]
cw = 2.95
for i, (title, color, items) in enumerate(scenarios):
    cx = 0.35 + i * (cw + 0.25)
    add_rect(s, cx, 1.4, cw, 0.55, fill_color=color)
    add_text(s, title, cx+0.1, 1.42, cw-0.2, 0.51, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, cx, 1.95, cw, 4.0, fill_color=WHITE, line_color=color, line_width=1)
    add_rect(s, cx, 1.95, cw, 0.4, fill_color=LIGHT_GRAY)
    add_text(s, "关键特性", cx+0.1, 1.98, cw-0.2, 0.32, font_size=13, bold=True, color=GRAY)
    y = 2.42
    for item in items:
        dot = s.shapes.add_shape(9, Inches(cx+0.15), Inches(y+0.12), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
        add_text(s, item, cx+0.32, y, cw-0.42, 0.6, font_size=14, color=DARK)
        y += 0.72
    add_rect(s, cx, 5.95, cw, 0.45, fill_color=LIGHT_GRAY)
    add_text(s, "截图区", cx, 6.05, cw, 0.32, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)
page_num(s, 7)

# ===== Slide 8: 技术架构 =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=BG)
add_title_bar(s, "技术架构", "Flutter Web + FastAPI + DeepSeek-V3 + 讯飞 + MediaPipe")
# 前端列
add_rect(s, 0.4, 1.4, 3.7, 0.5, fill_color=ORANGE)
add_text(s, "Flutter Web 前端", 0.5, 1.44, 3.5, 0.42, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, 0.4, 1.9, 3.7, 2.5, fill_color=RGBColor(0xFF,0xEE,0xE0))
front_items = ["长辈版 UI（17页面）", "AgentElementRegistry", "AgentFab 悬浮助手", "MediaPipe 本地人脸", "讯飞 WebSDK 语音"]
y = 2.0
for item in front_items:
    add_text(s, "• " + item, 0.55, y, 3.4, 0.42, font_size=14, color=DARK)
    y += 0.46
# 箭头
add_text(s, "⟺", 4.15, 2.55, 0.7, 0.5, font_size=24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, "WebSocket", 4.05, 3.05, 0.9, 0.3, font_size=11, color=GRAY, align=PP_ALIGN.CENTER)
# 后端列
add_rect(s, 4.95, 1.4, 3.7, 0.5, fill_color=GREEN)
add_text(s, "FastAPI 后端", 5.05, 1.44, 3.5, 0.42, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, 4.95, 1.9, 3.7, 2.5, fill_color=RGBColor(0xE8,0xF5,0xE9))
back_items = ["Agno Agent 框架", "DeepSeek-V3 推理", "WebSocket 状态机", "三级权限矩阵", "草稿箱管理"]
y = 2.0
for item in back_items:
    add_text(s, "• " + item, 5.1, y, 3.4, 0.42, font_size=14, color=DARK)
    y += 0.46
# 箭头
add_text(s, "→", 8.7, 2.55, 0.5, 0.5, font_size=24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
# AI服务列
ai_services = [("DeepSeek-V3", "意图理解 + 推理"), ("讯飞 ASR+TTS", "语音识别+合成"), ("MediaPipe", "本地人脸检测")]
y = 1.4
for name, desc in ai_services:
    add_rect(s, 9.3, y, 3.7, 0.45, fill_color=ORANGE)
    add_text(s, name, 9.4, y+0.05, 3.5, 0.35, font_size=14, bold=True, color=WHITE)
    add_rect(s, 9.3, y+0.45, 3.7, 0.4, fill_color=LIGHT_ORANGE)
    add_text(s, desc, 9.4, y+0.5, 3.5, 0.32, font_size=13, color=DARK)
    y += 1.0
# 工程亮点
add_rect(s, 0.4, 4.65, 12.5, 0.42, fill_color=ORANGE)
add_text(s, "工程安全亮点", 0.6, 4.7, 12.1, 0.32, font_size=15, bold=True, color=WHITE)
eng_items = [
    "确定性操作物理隔离（代码层永不注册为工具，提示词注入无效）",
    "权限状态随 WebSocket 会话消亡，不持久化，不泄漏",
    "草稿箱 IndexedDB 本地持久化，中途可续填，随时可撤销",
]
y = 5.15
for item in eng_items:
    add_text(s, "• " + item, 0.6, y, 12.1, 0.38, font_size=14, color=DARK)
    y += 0.42
page_num(s, 8)

# ===== Slide 9: 成果总结 =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=BG)
add_title_bar(s, "成果总结与学术贡献", "三项创新贡献")
contributions = [
    ("贡献一：设计框架", "受控响应型代理框架", [
        "三原则固化为架构约束",
        "L1/L2/L3 分级授权矩阵",
        "信任保障从对话约定\n升级为可验证工程性质",
    ], ORANGE),
    ("贡献二：方法论", "多模态冗余互补方法论", [
        "语音·视觉·触控三模态\n时序对齐",
        "任一感官退化不影响使用",
        "适老化设计的工程化路径",
    ], GREEN),
    ("贡献三：范式路径", "协助服务范式", [
        "政务大厅协助模式\n数字化复刻",
        "可推广至医疗/金融/生活",
        "协助而非替代，透明不黑箱",
    ], BLUE),
]
cw = 3.9
for i, (tag, title, items, color) in enumerate(contributions):
    cx = 0.4 + i * (cw + 0.3)
    add_rect(s, cx, 1.4, cw, 0.42, fill_color=color)
    add_text(s, tag, cx+0.1, 1.44, cw-0.2, 0.34, font_size=14, bold=True, color=WHITE)
    add_rect(s, cx, 1.82, cw, 0.52, fill_color=RGBColor(0xFF,0xEE,0xCC))
    add_text(s, title, cx+0.1, 1.86, cw-0.2, 0.44, font_size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_rect(s, cx, 2.34, cw, 3.5, fill_color=WHITE, line_color=color, line_width=1)
    y = 2.44
    for item in items:
        dot = s.shapes.add_shape(9, Inches(cx+0.15), Inches(y+0.12), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
        add_text(s, item, cx+0.32, y, cw-0.42, 0.8, font_size=15, color=DARK)
        y += 0.9
add_rect(s, 0.4, 6.05, 12.5, 0.55, fill_color=RGBColor(0xEE,0xEE,0xEE))
add_text(s, "可用性测试计划：5-8位老年用户 · 4场景任务 · 任务完成率 + SUS量表 + 信任感评分", 0.6, 6.12, 12.1, 0.4, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
page_num(s, 9)

# ===== Slide 10: 结尾 =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=ORANGE)
for r, c in [(6.0, RGBColor(0xFF,0x8F,0x00)), (4.0, RGBColor(0xFF,0xA0,0x30)), (2.5, RGBColor(0xFF,0xB0,0x50))]:
    shape = s.shapes.add_shape(9, Inches(9.5 - r/2), Inches(4.5 - r/2), Inches(r), Inches(r))
    shape.fill.solid(); shape.fill.fore_color.rgb = c; shape.line.fill.background()
add_text(s, "谢谢", 1.2, 1.0, 7.0, 1.5, font_size=80, bold=True, color=WHITE)
add_text(s, "代理是助手，用户始终是主角", 1.2, 2.8, 10.0, 0.7, font_size=28, bold=True, color=WHITE)
add_text(s, "协助而非替代，透明而非黑箱", 1.2, 3.55, 10.0, 0.7, font_size=28, bold=True, color=WHITE)
add_rect(s, 1.2, 4.5, 6.0, 0.05, fill_color=LIGHT_ORANGE)
add_text(s, "欢迎现场演示原型 · http://localhost:3080", 1.2, 4.7, 10.0, 0.5, font_size=18, color=LIGHT_ORANGE)
add_text(s, "信息服务APP的适老化设计与多模态交互", 1.2, 5.3, 10.0, 0.5, font_size=16, color=LIGHT_ORANGE)

output = "D:/Code/bs/docs/成果展示PPT.pptx"
prs.save(output)
print(f"saved: {output}")
